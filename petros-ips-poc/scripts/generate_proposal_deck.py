#!/usr/bin/env python3
"""
ABeam T260002 Proposal Deck Generator (v2 — full iteration set)
─────────────────────────────────────────────────────────────────
Reads PROPOSAL.md slide-by-slide, strips placeholder slides 36-84
from the existing template, and inserts production-grade replacements
in their place.

Iterations applied:
  1. POC screenshots auto-inserted on slides 60-71 from
     proposal-screenshots/ (captured by capture_proposal_screenshots.mjs)
  2. Native PowerPoint architecture diagrams on slides 37, 40, 76
  3. Status-pill shapes on the coverage scorecard (slide 52)
  4. Strip existing slides 36-84 before inserting new ones (clean output)
  5. Richer section dividers with big section number + agenda preview

Re-runnable. Source-of-truth: petros-ips-poc/PROPOSAL.md.
"""

import re
import sys
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ────────────────────────────────────────────────────────────
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
TEMPLATE = ROOT / "ABeam_PETROS_T260002_Final  -  Repaired.pptx"
PROPOSAL_MD = ROOT / "PROPOSAL.md"
SCREENSHOTS_DIR = ROOT / "proposal-screenshots"
OUTPUT = ROOT / "ABeam_T260002_Proposal_Generated.pptx"

# ── Brand colours (from theme1.xml inspection) ───────────────────────
NAVY = RGBColor(0x00, 0x19, 0x64)
TAUPE = RGBColor(0x7D, 0x6E, 0x59)
TAUPE_LIGHT = RGBColor(0xC8, 0xBE, 0xAA)
TAUPE_BG = RGBColor(0xF0, 0xEB, 0xE3)
BLUE_MID = RGBColor(0x45, 0x89, 0xAF)
BLUE_LIGHT = RGBColor(0xD9, 0xE5, 0xEC)
BLUE_PALE = RGBColor(0xE9, 0xF0, 0xF4)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xD4, 0xA8, 0x43)
GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Layout indices in the existing template ─────────────────────────
LAYOUT_FRONT_COVER = 0
LAYOUT_TITLE_CONTENT_JP = 1
LAYOUT_SECTION_DIVIDER = 2
LAYOUT_BACK_COVER = 3
LAYOUT_TITLE_ONLY_BLANK = 4
LAYOUT_TITLE_TEXT = 5
LAYOUT_TITLE_CONTENT = 6
LAYOUT_CONTENTS1 = 7
LAYOUT_TITLE_ONLY = 8
LAYOUT_TITLE_AND_CONTENTS = 12

# ── Slide-special routing ────────────────────────────────────────────
DIVIDER_SLIDES = {36, 43, 53, 60, 72, 77}
SECTION_NUMBERS = {
    36: ("§4", "Target Architecture", "SAC + Datasphere · architecture · integration · POC translation"),
    43: ("§5", "Functional Coverage", "RFP §1–§11 independent scoring · 503/503 tests · evidence per clause"),
    53: ("§6", "Differentiators", "62 PETROS-Sarawak deltas · two-pass methodology · why us"),
    60: ("§7", "POC Walkthrough", "Live system · 17 production pages · screenshots from petros-ips-poc.vercel.app"),
    72: ("§8", "Integration Design", "Datasphere connectors · 12 source systems · cross-tenant bridge"),
    77: ("§9", "Methodology & Phase Plan", "ABeam delivery · Phase 1a · Phase 1b · Phase 2 · hypercare"),
}
ARCHITECTURE_SLIDES = {37, 40, 76}
COVERAGE_SCORECARD_SLIDES = {52}
SCREENSHOT_SLIDES = {61, 62, 63, 64, 65, 66, 67, 68, 69, 70, 71}
SCREENSHOT_MAP = {
    61: "slide61-dashboard.png",
    62: "slide62-economics.png",
    63: "slide63-sensitivity.png",
    64: "slide64-monte-carlo.png",
    65: "slide65-portfolio.png",
    66: "slide66-financial.png",
    67: "slide67-reserves.png",
    68: "slide68-consolidation.png",
    69: "slide69-ma.png",
    70: "slide70-project-finance.png",
    71: "slide71-climate.png",
}


# ── Markdown parsing ─────────────────────────────────────────────────

def parse_proposal_md(path: Path):
    text = path.read_text(encoding="utf-8")
    slide_pattern = re.compile(r"^## Slide (\d+)\s*[—–-]\s*(.+?)$", re.MULTILINE)
    slides = []
    matches = list(slide_pattern.finditer(text))
    for i, m in enumerate(matches):
        num = int(m.group(1))
        title = m.group(2).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body_text = text[start:end].strip()
        slides.append({"num": num, "title": title, "raw": body_text})
    return slides


def parse_speaker_notes(raw):
    m = re.search(r"\*\*Speaker notes\*\*:\s*\n((?:>.*\n?)+)", raw)
    if not m:
        return ""
    notes = m.group(1)
    cleaned = re.sub(r"^>\s?", "", notes, flags=re.MULTILINE).strip()
    cleaned = cleaned.replace('"', '"').replace('"', '"').replace("'", "'").replace("'", "'")
    cleaned = re.sub(r'^"(.+)"$', r'\1', cleaned, flags=re.MULTILINE)
    return cleaned


def parse_body_md(raw):
    m = re.search(r"\*\*Speaker notes\*\*:", raw)
    body = raw[: m.start()] if m else raw
    body = re.sub(r"\*\*Visual notes\*\*[^\n]*\n", "", body)
    body = re.sub(r"\*\*Body\*\*[^\n]*\n+", "", body)
    return body.strip()


def strip_title_prefix(title):
    return re.sub(r"^Section divider:\s*", "", title).strip()


def split_body_into_blocks(body):
    blocks = []
    lines = body.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("###"):
            level = len(stripped) - len(stripped.lstrip("#"))
            text = stripped.lstrip("#").strip()
            blocks.append({"kind": "heading", "level": level, "text": text})
            i += 1
            continue
        if stripped.startswith("```"):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append({"kind": "code", "lines": code_lines})
            continue
        if stripped.startswith(">"):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(re.sub(r"^>\s?", "", lines[i]))
                i += 1
            blocks.append({"kind": "quote", "lines": quote_lines})
            continue
        if "|" in line and i + 1 < len(lines):
            sep = lines[i + 1].strip()
            if re.match(r"^\|?[\s:|-]+\|?$", sep) and "---" in sep:
                rows = [[c.strip() for c in line.strip().strip("|").split("|")]]
                i += 2
                while i < len(lines) and "|" in lines[i] and lines[i].strip():
                    rows.append([c.strip() for c in lines[i].strip().strip("|").split("|")])
                    i += 1
                blocks.append({"kind": "table", "rows": rows})
                continue
        bullet_match = re.match(r"^(\s*)[-*]\s+(.*)$", line)
        if bullet_match:
            bullets = []
            while i < len(lines):
                bm = re.match(r"^(\s*)[-*]\s+(.*)$", lines[i])
                if bm:
                    indent = len(bm.group(1)) // 2
                    text = bm.group(2).strip()
                    bullets.append({"level": indent, "text": text})
                    i += 1
                elif lines[i].strip() == "":
                    break
                elif lines[i].startswith("  ") and bullets:
                    bullets[-1]["text"] += " " + lines[i].strip()
                    i += 1
                else:
                    break
            blocks.append({"kind": "bullets", "items": bullets})
            continue
        num_match = re.match(r"^(\s*)\d+\.\s+(.*)$", line)
        if num_match:
            items = []
            while i < len(lines):
                nm = re.match(r"^(\s*)\d+\.\s+(.*)$", lines[i])
                if nm:
                    indent = len(nm.group(1)) // 2
                    text = nm.group(2).strip()
                    items.append({"level": indent, "text": text})
                    i += 1
                elif lines[i].strip() == "":
                    break
                else:
                    break
            blocks.append({"kind": "numbered", "items": items})
            continue
        para_lines = [line]
        i += 1
        while i < len(lines) and lines[i].strip() and not lines[i].strip().startswith(("#", ">", "-", "*", "```")) and "|" not in lines[i]:
            para_lines.append(lines[i])
            i += 1
        blocks.append({"kind": "para", "text": " ".join(p.strip() for p in para_lines)})
    return blocks


# ── Inline markdown rendering ────────────────────────────────────────

def render_runs_into_paragraph(p, text, base_size=11, base_color=TEXT_DARK):
    pos = 0
    tokens = []
    pattern = re.compile(r"(\*\*[^*]+?\*\*|`[^`]+`|\*[^*]+?\*)")
    for m in pattern.finditer(text):
        if m.start() > pos:
            tokens.append(("text", text[pos:m.start()]))
        seg = m.group(0)
        if seg.startswith("**"):
            tokens.append(("bold", seg[2:-2]))
        elif seg.startswith("`"):
            tokens.append(("code", seg[1:-1]))
        elif seg.startswith("*"):
            tokens.append(("italic", seg[1:-1]))
        pos = m.end()
    if pos < len(text):
        tokens.append(("text", text[pos:]))
    if not tokens:
        tokens = [("text", text)]
    for kind, t in tokens:
        run = p.add_run()
        run.text = t
        run.font.size = Pt(base_size)
        run.font.color.rgb = base_color
        if kind == "bold":
            run.font.bold = True
            run.font.color.rgb = NAVY
        elif kind == "code":
            run.font.name = "Consolas"
            run.font.size = Pt(base_size - 1)
            run.font.color.rgb = TAUPE
        elif kind == "italic":
            run.font.italic = True


# ─────────────────────────────────────────────────────────────────────
# ITERATION 4 — Strip existing slides 36-84 from template
# ─────────────────────────────────────────────────────────────────────

def strip_existing_36_to_84(prs):
    """Remove slides 36-84 (1-indexed) from the template before appending
    new ones.  Cleanly drops both the sldIdLst entry AND the underlying
    relationship, so the saved file has no orphan slide parts.

    Earlier versions of this function only removed the sldIdLst entry —
    that left 49 orphan slide XML files in the zip + 49 dangling rels in
    presentation.xml.rels.  Strict PowerPoint validators (and the file-
    repair dialog) flag that as content corruption."""
    rels_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    to_delete = [(i, e) for i, e in enumerate(sld_ids) if 35 <= i <= 83]
    if not to_delete:
        return 0
    for idx, elem in reversed(to_delete):
        rId = elem.attrib.get(f"{rels_ns}id")
        # Drop the relationship from presentation part; the orphan slide
        # part is then GC-eligible and python-pptx omits it on save.
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(elem)
    return len(to_delete)


def reorder_slides_to_canonical(prs, n_kept_at_start=35, n_new=49):
    """After strip + append, re-order the sldIdLst so the layout is:
       [original 1-35][new 36-84][original 85-98]

    After strip, sldIdLst contained: original_1..35  +  original_85..98
    (49 entries). After appending n_new (=49) new slides, sldIdLst now has
    98 entries arranged as either:
       (a) original_1..35 + original_85..98 + new_36..84   — what we want at the end
    or (b) some interleaving caused by python-pptx insertion behaviour.

    Either way, this function moves the n_new most-recently-appended
    children to land between index n_kept_at_start and n_kept_at_start + n_new,
    so the final ordering is:
       [0..n_kept_at_start-1]: original 1..35
       [n_kept_at_start..n_kept_at_start+n_new-1]: new 36..84
       [n_kept_at_start+n_new..]: original 85..98
    """
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    if len(children) < n_kept_at_start + n_new:
        return  # nothing sensible to do
    # Identify the new ones — they are the LAST n_new children (added last)
    head = children[:n_kept_at_start]
    middle_old = children[n_kept_at_start: len(children) - n_new]
    new_block = children[len(children) - n_new:]
    desired = head + new_block + middle_old
    # Detach all children, re-add in desired order
    for c in children:
        sldIdLst.remove(c)
    for c in desired:
        sldIdLst.append(c)


# ─────────────────────────────────────────────────────────────────────
# ITERATION 5 — Richer section dividers
# ─────────────────────────────────────────────────────────────────────

def build_section_divider_v2(prs, slide_num, title, blocks):
    """Richer section divider:
      - Big section number (§4, §5, ...) in taupe on left third
      - Section name (large, navy) on right two-thirds
      - Subtitle / agenda preview below
    Uses 中扉 layout as the base for the brand identity, then overlays."""
    layout = prs.slide_layouts[LAYOUT_SECTION_DIVIDER]
    slide = prs.slides.add_slide(layout)

    # Hide existing placeholders by clearing their text
    # (we'll overlay our own shapes for the rich divider)
    for ph in slide.placeholders:
        if hasattr(ph, "text_frame") and ph.text_frame is not None:
            ph.text = ""

    section_num, section_name, agenda_preview = SECTION_NUMBERS.get(
        slide_num, ("", strip_title_prefix(title), "")
    )

    # Big section number (left, large, taupe)
    # Slide is 720x405pt. Number occupies left ~30%.
    num_box = slide.shapes.add_textbox(Pt(36), Pt(120), Pt(180), Pt(180))
    num_tf = num_box.text_frame
    num_tf.margin_left = Pt(0)
    num_tf.margin_top = Pt(0)
    num_p = num_tf.paragraphs[0]
    num_p.alignment = PP_ALIGN.LEFT
    num_run = num_p.add_run()
    num_run.text = section_num
    num_run.font.name = "Arial"
    num_run.font.size = Pt(120)
    num_run.font.bold = True
    num_run.font.color.rgb = TAUPE_LIGHT  # subtle, decorative

    # Vertical divider line between number and content
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Pt(225), Pt(120), Pt(2), Pt(180)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TAUPE
    line.line.fill.background()

    # Section name (right, large, navy)
    name_box = slide.shapes.add_textbox(Pt(245), Pt(140), Pt(440), Pt(60))
    name_tf = name_box.text_frame
    name_tf.margin_left = Pt(0)
    name_tf.margin_top = Pt(0)
    name_p = name_tf.paragraphs[0]
    name_run = name_p.add_run()
    name_run.text = section_name
    name_run.font.name = "Arial"
    name_run.font.size = Pt(36)
    name_run.font.bold = True
    name_run.font.color.rgb = NAVY

    # Agenda preview / tagline (right, smaller, taupe)
    agenda_box = slide.shapes.add_textbox(Pt(245), Pt(210), Pt(440), Pt(80))
    agenda_tf = agenda_box.text_frame
    agenda_tf.word_wrap = True
    agenda_tf.margin_left = Pt(0)
    agenda_p = agenda_tf.paragraphs[0]
    agenda_run = agenda_p.add_run()
    agenda_run.text = agenda_preview
    agenda_run.font.name = "Arial"
    agenda_run.font.size = Pt(13)
    agenda_run.font.italic = True
    agenda_run.font.color.rgb = TAUPE

    return slide


# ─────────────────────────────────────────────────────────────────────
# ITERATION 3 — Status pills
# ─────────────────────────────────────────────────────────────────────

def render_status_pill(slide, left, top, width, height, status_token):
    """Native PowerPoint rounded rectangle with colour by status."""
    # Map status emoji/marker to colour + label
    mapping = {
        "✅": (GREEN, WHITE, "MET"),
        "✓": (GREEN, WHITE, "MET"),
        "◐": (AMBER, WHITE, "PARTIAL"),
        "⚠": (RED, WHITE, "GAP"),
        "⊘": (GRAY, WHITE, "OUT"),
    }
    fill_color, text_color, label = mapping.get(status_token, (GRAY, WHITE, status_token))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Arial"
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shape


# ─────────────────────────────────────────────────────────────────────
# ITERATION 2 — Architecture diagrams
# ─────────────────────────────────────────────────────────────────────

def add_arrow(slide, x1, y1, x2, y2, color=TAUPE, weight=2.0):
    """Draw a simple arrow between two points."""
    from pptx.shapes.connector import Connector
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)


def build_architecture_overview_slide(prs, slide_num, title, blocks, body_blocks):
    """Slide 37 — 3-layer architecture diagram + supporting table."""
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    # Title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = NAVY

    # 3-layer diagram in left half
    layer_x = Pt(36)
    layer_w = Pt(330)
    layer_h = Pt(70)
    layer_gap = Pt(15)
    layer_y_start = Pt(85)

    layers = [
        ("Planning & Analytics — SAP SAC", "Stories · DAS · Multi-Actions · Calculated Measures · Smart Predict", BLUE_PALE, NAVY),
        ("Data Fabric & Modelling — SAP Datasphere", "Federated query · semantic layer · cross-tenant bridge · master-data governance", BLUE_LIGHT, NAVY),
        ("Source systems — existing PETROS landscape", "S/4HANA · SuccessFactors · Ariba · DRC · DMS · External (Brent / LNG / FX)", TAUPE_BG, TAUPE),
    ]
    for i, (lbl, desc, fill, line_color) in enumerate(layers):
        y = Pt(85 + i * 85)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, layer_x, y, layer_w, layer_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = line_color
        rect.line.width = Pt(1.5)
        tf = rect.text_frame
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = lbl + "\n"
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = TEXT_DARK

    # Vertical arrows between layers (bottom-up)
    for i in range(2):
        y_top = Pt(85 + i * 85 + 70 + 2)
        y_bot = Pt(85 + (i + 1) * 85 - 2)
        cx = layer_x + layer_w // 2
        line = slide.shapes.add_connector(1, cx - Pt(20), y_top, cx - Pt(20), y_bot)
        line.line.color.rgb = TAUPE
        line.line.width = Pt(1.5)
        line2 = slide.shapes.add_connector(1, cx + Pt(20), y_bot, cx + Pt(20), y_top)
        line2.line.color.rgb = NAVY
        line2.line.width = Pt(1.5)

    # Caption
    cap_box = slide.shapes.add_textbox(layer_x, Pt(345), layer_w, Pt(20))
    cap = cap_box.text_frame.paragraphs[0]
    cap_run = cap.add_run()
    cap_run.text = "← actuals & master data    write-back of approved plans →"
    cap_run.font.size = Pt(8)
    cap_run.font.italic = True
    cap_run.font.color.rgb = TEXT_MUTED

    # Right side: short callout list
    right_x = Pt(385)
    right_y = Pt(85)
    callout_box = slide.shapes.add_textbox(right_x, right_y, Pt(300), Pt(280))
    ctf = callout_box.text_frame
    ctf.word_wrap = True
    callouts = [
        ("All three layers SAP-native", "No new infrastructure. Datasphere connects what PETROS already operates."),
        ("POC = executable functional spec", "Engine modules (engine/*) become Datasphere views and SAC DAS scripts."),
        ("Cross-tenant ready", "Datasphere bridges PETROS S/4 ↔ PETRONAS-Group landscape under transition (D9)."),
        ("Power BI as alternate consumer", "Same Datasphere views feed both SAC (planning) and Power BI (executive)."),
    ]
    first = True
    for hdr, body in callouts:
        p = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = "▸ "
        r.font.size = Pt(11)
        r.font.color.rgb = TAUPE
        r2 = p.add_run()
        r2.text = hdr + "\n"
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = NAVY
        r3 = p.add_run()
        r3.text = body
        r3.font.size = Pt(9)
        r3.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    return slide


def build_planning_cycle_slide(prs, slide_num, title, blocks):
    """Slide 40 — 6-step planning cycle as left-to-right chevron flow."""
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = NAVY

    # Subtitle
    sub_box = slide.shapes.add_textbox(Pt(36), Pt(85), Pt(648), Pt(20))
    sp = sub_box.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Continuous · auditable · write-back enabled"
    sr.font.size = Pt(11)
    sr.font.italic = True
    sr.font.color.rgb = TAUPE

    steps = [
        ("1", "S/4HANA", "ACDOCA / CSKS / PRPS\nactuals + master data", BLUE_PALE),
        ("2", "Datasphere", "Federated query\nno copy, no delay", BLUE_LIGHT),
        ("3", "SAC", "actuals version\nmonthly close", BLUE_MID),
        ("4", "Planner", "working version\nwhat-if + edits", TAUPE_BG),
        ("5", "Workflow", "Submit → Review → Approve\nSoD-enforced", TAUPE_LIGHT),
        ("6", "Write-back", "S/4 capex commit\nDRC tax filing", TAUPE),
    ]
    n = len(steps)
    avail_w = 648
    chev_w = avail_w / n
    chev_y = Pt(125)
    chev_h = Pt(85)

    for i, (num, lbl, desc, fill) in enumerate(steps):
        x = Pt(36 + i * chev_w)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON if i < n - 1 else MSO_SHAPE.ROUNDED_RECTANGLE,
            x, chev_y, Pt(chev_w - 4), chev_h,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = NAVY if i < 3 else TAUPE
        rect.line.width = Pt(1.0)
        tf = rect.text_frame
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = num + "  "
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = NAVY if i < 3 else TAUPE
        r2 = p.add_run()
        r2.text = lbl + "\n"
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = NAVY
        r3 = p.add_run()
        r3.text = desc
        r3.font.size = Pt(8)
        r3.font.color.rgb = TEXT_DARK

    # Cadence strip below
    cad_box = slide.shapes.add_textbox(Pt(36), Pt(230), Pt(648), Pt(28))
    ctf = cad_box.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    runs = [
        ("monthly close ", BLUE_MID, True),
        ("→  ", TAUPE, False),
        ("quarterly forecast ", AMBER, True),
        ("→  ", TAUPE, False),
        ("annual budget freeze", NAVY, True),
    ]
    for txt, col, bold in runs:
        r = cp.add_run()
        r.text = txt
        r.font.size = Pt(10)
        r.font.bold = bold
        r.font.color.rgb = col

    # Step-detail table below
    detail_table_top = Pt(265)
    table_data = [
        ["Step", "Where", "What", "Cadence"],
        ["1", "S/4HANA ACDOCA / CSKS / PRPS", "Actuals + master data published", "continuous"],
        ["2", "Datasphere", "Federated query — no copy, no delay", "continuous"],
        ["3", "SAC Planning model (`actuals` version)", "Loaded into 6-version registry", "monthly"],
        ["4", "SAC Story (`working` → `submitted`)", "Planner edits + what-if analysis", "per-cycle"],
        ["5", "SAC Calendar workflow", "SoD-enforced submit → review → approve", "per-cycle"],
        ["6", "DRC + S/4 write-back", "PITA / SST filings + CAPEX commitments", "post-approval"],
    ]
    n_rows = len(table_data)
    n_cols = len(table_data[0])
    tbl_left = Pt(36)
    tbl_width = Pt(648)
    tbl_height = Pt(min(20 * n_rows, 130))
    gfx = slide.shapes.add_table(n_rows, n_cols, tbl_left, detail_table_top, tbl_width, tbl_height)
    tbl = gfx.table
    for ci, hdr in enumerate(table_data[0]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        ctf2 = cell.text_frame
        ctf2.clear()
        cp2 = ctf2.paragraphs[0]
        cr = cp2.add_run()
        cr.text = hdr
        cr.font.size = Pt(8)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
    for ri in range(1, n_rows):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TAUPE_BG if ri % 2 == 1 else WHITE
            ctf2 = cell.text_frame
            ctf2.clear()
            cp2 = ctf2.paragraphs[0]
            cr = cp2.add_run()
            cr.text = table_data[ri][ci].replace("`", "")
            cr.font.size = Pt(8)
            cr.font.color.rgb = TEXT_DARK

    return slide


def build_cross_tenant_bridge_slide(prs, slide_num, title, blocks):
    """Slide 76 — cross-tenant Datasphere bridge diagram."""
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = NAVY

    # Two tenants on top row
    tenant_y = Pt(95)
    tenant_w = Pt(220)
    tenant_h = Pt(95)

    # Left: PETRONAS-Group S/4
    left_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(60), tenant_y, tenant_w, tenant_h)
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = BLUE_PALE
    left_rect.line.color.rgb = NAVY
    left_rect.line.width = Pt(1.5)
    ltf = left_rect.text_frame
    ltf.margin_left = Pt(10)
    ltf.margin_top = Pt(6)
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    r1 = lp.add_run()
    r1.text = "PETRONAS-Group S/4HANA tenant\n"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r2 = lp.add_run()
    r2.text = "Existing landscape (slide 35)\nGL · CSKS · PRPS · ACDOCA\nJVA · Procurement · Asset Mgmt"
    r2.font.size = Pt(9)
    r2.font.color.rgb = TEXT_DARK

    # Right: PETROS-dedicated S/4 (future state, dashed border)
    right_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(440), tenant_y, tenant_w, tenant_h)
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = TAUPE_BG
    right_rect.line.color.rgb = TAUPE
    right_rect.line.width = Pt(1.5)
    # Note: dashed border isn't trivial via python-pptx; flag in subtitle
    rtf = right_rect.text_frame
    rtf.margin_left = Pt(10)
    rtf.margin_top = Pt(6)
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rr1 = rp.add_run()
    rr1.text = "PETROS-dedicated tenant (future)\n"
    rr1.font.size = Pt(12)
    rr1.font.bold = True
    rr1.font.color.rgb = TAUPE
    rr2 = rp.add_run()
    rr2.text = "Plausible post-CSA evolution\n(not yet committed)\nBridge accommodates either path"
    rr2.font.size = Pt(9)
    rr2.font.italic = True
    rr2.font.color.rgb = TEXT_DARK

    # Datasphere bridge in the middle (centered, prominent)
    bridge_y = Pt(220)
    bridge_w = Pt(440)
    bridge_h = Pt(60)
    bridge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(150), bridge_y, bridge_w, bridge_h)
    bridge.fill.solid()
    bridge.fill.fore_color.rgb = NAVY
    bridge.line.fill.background()
    btf = bridge.text_frame
    btf.margin_left = Pt(10)
    btf.margin_top = Pt(6)
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br1 = bp.add_run()
    br1.text = "SAP Datasphere — cross-tenant federation\n"
    br1.font.size = Pt(13)
    br1.font.bold = True
    br1.font.color.rgb = WHITE
    br2 = bp.add_run()
    br2.text = "Federated query · row-level masking · view-level access · audit log on every query"
    br2.font.size = Pt(9)
    br2.font.color.rgb = WHITE

    # Connectors from each tenant down to the bridge
    for tx in [Pt(170), Pt(550)]:
        line = slide.shapes.add_connector(1, tx, tenant_y + tenant_h, tx, bridge_y)
        line.line.color.rgb = TAUPE
        line.line.width = Pt(1.5)

    # SAC consumes from bridge — bottom box
    sac_y = Pt(305)
    sac = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(280), sac_y, Pt(180), Pt(50))
    sac.fill.solid()
    sac.fill.fore_color.rgb = BLUE_MID
    sac.line.fill.background()
    stf = sac.text_frame
    stf.margin_left = Pt(8)
    stf.margin_top = Pt(4)
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "SAC IPS Tenant\n"
    sr.font.size = Pt(12)
    sr.font.bold = True
    sr.font.color.rgb = WHITE
    sr2 = sp.add_run()
    sr2.text = "Planning · workflow · disclosure"
    sr2.font.size = Pt(9)
    sr2.font.color.rgb = WHITE

    # Connector from bridge down to SAC
    conn = slide.shapes.add_connector(1, Pt(370), bridge_y + bridge_h, Pt(370), sac_y)
    conn.line.color.rgb = NAVY
    conn.line.width = Pt(2.0)

    # Caption strip
    cap_box = slide.shapes.add_textbox(Pt(36), Pt(365), Pt(648), Pt(25))
    cp = cap_box.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr1 = cp.add_run()
    cr1.text = "D9 closure: cross-tenant Datasphere federation per SAP-published architecture pattern. No tenant migration required."
    cr1.font.size = Pt(9)
    cr1.font.italic = True
    cr1.font.color.rgb = TEXT_MUTED

    return slide


# ─────────────────────────────────────────────────────────────────────
# ITERATION 1 — POC screenshot insertion
# ─────────────────────────────────────────────────────────────────────

def build_screenshot_slide(prs, slide_num, title, blocks):
    """POC walkthrough slide (60-71) — title + screenshot + speaker callouts."""
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = NAVY

    # Insert screenshot if available
    fname = SCREENSHOT_MAP.get(slide_num)
    sshot_path = SCREENSHOTS_DIR / fname if fname else None
    img_left = Pt(36)
    img_top = Pt(80)
    img_width = Pt(440)
    img_height = Pt(248)

    if sshot_path and sshot_path.exists():
        try:
            slide.shapes.add_picture(
                str(sshot_path), img_left, img_top, width=img_width, height=img_height,
            )
            # Add a thin border by overlaying a transparent rectangle
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height
            )
            border.fill.background()
            border.line.color.rgb = TAUPE
            border.line.width = Pt(0.75)
        except Exception as e:
            print(f"  Warning: couldn't insert screenshot for slide {slide_num}: {e}")
    else:
        # Placeholder rectangle
        ph_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height
        )
        ph_rect.fill.solid()
        ph_rect.fill.fore_color.rgb = TAUPE_BG
        ph_rect.line.color.rgb = TAUPE
        ph_rect.line.width = Pt(0.75)
        ptf = ph_rect.text_frame
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = f"[POC screenshot — {fname or 'capture pending'}]"
        pr.font.size = Pt(11)
        pr.font.italic = True
        pr.font.color.rgb = TEXT_MUTED

    # Right column: condensed body text (key bullets only)
    right_x = Pt(490)
    right_y = Pt(80)
    right_w = Pt(195)
    right_h = Pt(280)
    right_box = slide.shapes.add_textbox(right_x, right_y, right_w, right_h)
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Pt(0)
    first = True
    # Take first paragraph as headline subtitle, then any bullets/headings
    headline = next((b["text"] for b in blocks if b["kind"] == "para"), "")
    if headline:
        p = rtf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = headline.split("—")[0].strip()
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = NAVY
        if "—" in headline:
            r2 = p.add_run()
            r2.text = "  " + headline.split("—", 1)[1].strip()
            r2.font.size = Pt(9)
            r2.font.italic = True
            r2.font.color.rgb = TAUPE
        p.space_after = Pt(8)

    # Render headings + bullets in compact form
    for block in blocks:
        if block["kind"] == "heading":
            p = rtf.paragraphs[0] if first else rtf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = block["text"]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_after = Pt(2)
        elif block["kind"] == "bullets":
            for item in block["items"]:
                p = rtf.paragraphs[0] if first else rtf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = "• "
                run.font.size = Pt(9)
                run.font.color.rgb = TAUPE
                render_runs_into_paragraph(p, item["text"], base_size=8)
                p.space_after = Pt(2)
    return slide


# ─────────────────────────────────────────────────────────────────────
# Standard content slide (existing behaviour)
# ─────────────────────────────────────────────────────────────────────

def build_content_slide(prs, slide_num, title, blocks):
    layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = NAVY

    body_left = Pt(36)
    body_top = Pt(78)
    body_width = Pt(720 - 72)
    body_height = Pt(405 - 78 - 25)
    content_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    first = True

    # Determine if this is the coverage scorecard slide (special status-pill rendering)
    is_scorecard = slide_num in COVERAGE_SCORECARD_SLIDES

    for block in blocks:
        if block["kind"] == "para":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            render_runs_into_paragraph(p, block["text"], base_size=11)
            p.space_after = Pt(6)
        elif block["kind"] == "heading":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = block["text"]
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_after = Pt(4)
        elif block["kind"] == "bullets":
            for item in block["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.level = min(item["level"], 4)
                bullet_run = p.add_run()
                bullet_run.text = "•  "
                bullet_run.font.size = Pt(11)
                bullet_run.font.color.rgb = TAUPE
                render_runs_into_paragraph(p, item["text"], base_size=11)
                p.space_after = Pt(2)
        elif block["kind"] == "numbered":
            for n, item in enumerate(block["items"], start=1):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.level = min(item["level"], 4)
                bullet_run = p.add_run()
                bullet_run.text = f"{n}.  "
                bullet_run.font.size = Pt(11)
                bullet_run.font.bold = True
                bullet_run.font.color.rgb = NAVY
                render_runs_into_paragraph(p, item["text"], base_size=11)
                p.space_after = Pt(2)
        elif block["kind"] == "code":
            for line in block["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = line
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                run.font.color.rgb = TAUPE
                p.space_after = Pt(0)
        elif block["kind"] == "quote":
            for line in block["lines"]:
                if not line.strip():
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = line.strip()
                run.font.size = Pt(11)
                run.font.italic = True
                run.font.color.rgb = TAUPE
                p.space_after = Pt(3)

    # Tables
    table_blocks = [b for b in blocks if b["kind"] == "table"]
    if table_blocks:
        from pptx.util import Inches
        non_table_count = sum(1 for b in blocks if b["kind"] != "table")
        table_top = Inches(2.4) if non_table_count > 0 else Inches(1.05)
        for ti, tbl_block in enumerate(table_blocks):
            rows = tbl_block["rows"]
            n_rows = len(rows)
            n_cols = len(rows[0]) if rows else 0
            if n_rows < 2 or n_cols < 1:
                continue
            tbl_left = Inches(0.4)
            tbl_width = Inches(9.2)
            tbl_height = Inches(min(0.28 * n_rows + 0.05, 3.5))
            try:
                gfx = slide.shapes.add_table(n_rows, n_cols, tbl_left, table_top, tbl_width, tbl_height)
                table = gfx.table
                # Header row
                for ci, cell_text in enumerate(rows[0]):
                    cell = table.cell(0, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                    tfc = cell.text_frame
                    tfc.clear()
                    p = tfc.paragraphs[0]
                    run = p.add_run()
                    run.text = cell_text.replace("**", "")
                    run.font.size = Pt(8)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                # Body rows
                for ri in range(1, n_rows):
                    row_data = rows[ri]
                    for ci in range(n_cols):
                        cell = table.cell(ri, min(ci, len(row_data) - 1)) if row_data else table.cell(ri, ci)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = TAUPE_BG if ri % 2 == 1 else WHITE
                        tfc = cell.text_frame
                        tfc.clear()
                        p = tfc.paragraphs[0]
                        cell_text = row_data[ci] if ci < len(row_data) else ""
                        clean = cell_text.replace("**", "").replace("`", "")

                        # Status-pill rendering for scorecard slide (D52)
                        if is_scorecard and ci == 1:  # 'Status' is column 1 in slide 52
                            for token in ("✅", "◐", "⚠", "⊘", "✓"):
                                if token in clean:
                                    # Render pill in cell
                                    cw = table.columns[ci].width
                                    rh = table.rows[ri].height
                                    # We can't easily place a shape inside a cell; instead
                                    # set the cell text to the colored label.
                                    p.alignment = PP_ALIGN.CENTER
                                    color_map = {
                                        "✅": (GREEN, "MET"),
                                        "✓": (GREEN, "MET"),
                                        "◐": (AMBER, "PARTIAL"),
                                        "⚠": (RED, "GAP"),
                                        "⊘": (GRAY, "OUT"),
                                    }
                                    pill_color, pill_label = color_map[token]
                                    # Set cell fill to pill colour, text white
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = pill_color
                                    run = p.add_run()
                                    run.text = pill_label
                                    run.font.size = Pt(7)
                                    run.font.bold = True
                                    run.font.color.rgb = WHITE
                                    break
                            else:
                                run = p.add_run()
                                run.text = clean
                                run.font.size = Pt(8)
                                run.font.color.rgb = TEXT_DARK
                        else:
                            run = p.add_run()
                            run.text = clean
                            run.font.size = Pt(8)
                            run.font.color.rgb = TEXT_DARK
                table_top = Inches(table_top.inches + tbl_height.inches + 0.08)
            except Exception as e:
                print(f"  Warning: failed to render table on slide {slide_num}: {e}")
                continue

    return slide


# ─────────────────────────────────────────────────────────────────────
# Speaker notes + footer
# ─────────────────────────────────────────────────────────────────────

def add_speaker_notes(slide, notes_text):
    if not notes_text:
        return
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return
    cleaned = notes_text.strip().replace('"', "").replace('"', "").replace('"', "")
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", cleaned) if p.strip()]
    if not paragraphs:
        return
    body_ph = None
    for ph in notes_slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph is not None and hasattr(body_ph, "text_frame") and body_ph.text_frame is not None:
        tf = body_ph.text_frame
        tf.text = paragraphs[0]
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
        for extra in paragraphs[1:]:
            new_p = tf.add_paragraph()
            run = new_p.add_run()
            run.text = extra
            run.font.size = Pt(10)
        return
    try:
        ntf = notes_slide.notes_text_frame
        if ntf is not None:
            ntf.text = "\n\n".join(paragraphs)
            return
    except Exception:
        pass


def add_footer(slide, slide_num, section_label):
    left = Pt(0)
    top = Pt(390)
    width = Pt(720)
    height = Pt(15)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"{slide_num}    "
    r1.font.size = Pt(8)
    r1.font.color.rgb = TEXT_MUTED
    r2 = p.add_run()
    r2.text = "TENDER NO: T260002 Integrated Planning System Implementation Proposal"
    r2.font.size = Pt(8)
    r2.font.color.rgb = TEXT_MUTED
    r3 = p.add_run()
    r3.text = "          confidential   ©2026 ABeam Consulting (Malaysia)"
    r3.font.size = Pt(8)
    r3.font.color.rgb = TEXT_MUTED


def section_label_for_slide(num):
    if 36 <= num <= 42: return "§4 Target Architecture"
    if 43 <= num <= 52: return "§5 Functional Coverage"
    if 53 <= num <= 59: return "§6 PETROS-Sarawak Differentiators"
    if 60 <= num <= 71: return "§7 POC Walkthrough"
    if 72 <= num <= 76: return "§8 Integration Design"
    if 77 <= num <= 84: return "§9 Methodology & Phase Plan"
    return ""


# ─────────────────────────────────────────────────────────────────────
# Main pipeline
# ─────────────────────────────────────────────────────────────────────

def main():
    print(f"→ Reading template: {TEMPLATE.name}")
    if not TEMPLATE.exists():
        print(f"  ERROR: template not found at {TEMPLATE}")
        sys.exit(1)
    print(f"→ Reading proposal: {PROPOSAL_MD.name}")
    if not PROPOSAL_MD.exists():
        print(f"  ERROR: PROPOSAL.md not found at {PROPOSAL_MD}")
        sys.exit(1)

    slides_data = parse_proposal_md(PROPOSAL_MD)
    print(f"  parsed {len(slides_data)} slide sections from PROPOSAL.md")

    prs = Presentation(str(TEMPLATE))
    print(f"  template has {len(prs.slides)} existing slides + {len(prs.slide_layouts)} layouts")

    # ── Iteration 4: strip existing slides 36-84 ─────────────────────
    deleted = strip_existing_36_to_84(prs)
    print(f"→ Stripped {deleted} existing slides at positions 36-84")
    print(f"  template now has {len(prs.slides)} slides remaining")

    # Check screenshot availability
    screenshots_available = SCREENSHOTS_DIR.exists() and any(SCREENSHOTS_DIR.glob("*.png"))
    print(f"→ Screenshots dir exists: {SCREENSHOTS_DIR.exists()}; files: {len(list(SCREENSHOTS_DIR.glob('*.png'))) if SCREENSHOTS_DIR.exists() else 0}")

    print(f"→ Generating slides 36-84 (with iteration overlays)")
    built = 0
    for s in slides_data:
        if s["num"] < 36 or s["num"] > 84:
            continue
        body = parse_body_md(s["raw"])
        notes = parse_speaker_notes(s["raw"])
        blocks = split_body_into_blocks(body)

        try:
            num = s["num"]
            if num in DIVIDER_SLIDES:
                slide = build_section_divider_v2(prs, num, s["title"], blocks)
            elif num == 37:
                slide = build_architecture_overview_slide(prs, num, s["title"], blocks, blocks)
            elif num == 40:
                slide = build_planning_cycle_slide(prs, num, s["title"], blocks)
            elif num == 76:
                slide = build_cross_tenant_bridge_slide(prs, num, s["title"], blocks)
            elif num in SCREENSHOT_SLIDES:
                slide = build_screenshot_slide(prs, num, s["title"], blocks)
            else:
                slide = build_content_slide(prs, num, s["title"], blocks)

            add_speaker_notes(slide, notes)
            add_footer(slide, num, section_label_for_slide(num))
            built += 1
            if built % 10 == 0:
                print(f"  built {built}…")
        except Exception as e:
            print(f"  WARN slide {s['num']!r}: {e}")
            import traceback; traceback.print_exc()

    print(f"  built {built} slides total")

    # Reorder so [original 1-35][new 36-84][original 85-98]
    reorder_slides_to_canonical(prs, n_kept_at_start=35, n_new=built)
    print(f"→ Re-ordered slides to canonical layout (1-35 | new 36-84 | 85+)")

    print(f"→ Saving to {OUTPUT.name}")
    prs.save(str(OUTPUT))
    print(f"  ✓ saved ({OUTPUT.stat().st_size // 1024}KB)")
    print(f"  total slides in output: {len(prs.slides)}")


if __name__ == "__main__":
    main()
